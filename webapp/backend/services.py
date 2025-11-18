from __future__ import annotations

import asyncio
import json
import os
import re
import uuid
from pathlib import Path
from tempfile import NamedTemporaryFile
from typing import Any, AsyncGenerator, Dict, List, Optional, Tuple

from litellm import acompletion, completion
from pptx import Presentation
from pypdf import PdfReader

from server.db import get_conn
from server.tools.canvas_export import render_canvas_markdown
from .schemas import (
    PaperChatMessage,
    Question,
    QuestionContextUploadResponse,
    QuestionGenerationRequest,
    QuestionGenerationResponse,
)

# DEFAULT_MODEL = os.getenv("LITELLM_MODEL", "gpt-4o-mini")
DEFAULT_MODEL = os.getenv("LITELLM_MODEL", "gpt-5-mini")
_raw_temp = os.getenv("LITELLM_TEMPERATURE")
if _raw_temp is not None:
    TEMPERATURE = float(_raw_temp)
elif "gpt-5" in DEFAULT_MODEL.lower():
    TEMPERATURE = 1.0
else:
    TEMPERATURE = 0.2
MAX_TOKENS = int(os.getenv("LITELLM_MAX_TOKENS", "4000"))
MAX_CONTEXT_CHARS = int(os.getenv("QUESTION_CONTEXT_CHAR_LIMIT", "60000"))
TYPE_PATTERNS = {
    "mcq": r"mcqs?|multiple\s+choice(?:\s+questions?)?",
    "short_answer": r"short[-\s]?answer(?:\s+questions?)?",
    "true_false": r"true\s*(?:or|/)?\s*false(?:\s+questions?)?|tf",
    "essay": r"essay(?:\s+questions?)?",
}


class QuestionGenerationError(RuntimeError):
    """Raised when the LLM output cannot be parsed."""


def _completion_limit_args(model_name: str) -> Dict[str, Any]:
    """Return the correct token-limit argument for the selected model."""
    if "gpt-5" in (model_name or "").lower():
        return {"max_completion_tokens": MAX_TOKENS}
    return {"max_tokens": MAX_TOKENS}


def generate_questions(payload: QuestionGenerationRequest) -> QuestionGenerationResponse:
    if not os.getenv("OPENAI_API_KEY") and not os.getenv("LITELLM_API_KEY"):
        raise QuestionGenerationError("OPENAI_API_KEY (or LITELLM_API_KEY) must be set to use the question generator.")

    messages = _build_messages(payload)
    comp_kwargs = _completion_limit_args(DEFAULT_MODEL)
    try:
        response = completion(
            model=DEFAULT_MODEL,
            messages=messages,
            temperature=TEMPERATURE,
            **comp_kwargs,
        )
    except Exception as exc:
        raise QuestionGenerationError(f"LLM request failed: {exc}") from exc

    content = response["choices"][0]["message"]["content"]
    questions = _parse_questions(content)
    markdown = render_canvas_markdown(payload.instructions, [q.model_dump() for q in questions], {})

    return QuestionGenerationResponse(
        questions=questions,
        markdown=markdown,
        raw_response=content,
    )


async def stream_generate_questions(payload: QuestionGenerationRequest) -> AsyncGenerator[Dict[str, Any], None]:
    if not os.getenv("OPENAI_API_KEY") and not os.getenv("LITELLM_API_KEY"):
        raise QuestionGenerationError("OPENAI_API_KEY (or LITELLM_API_KEY) must be set to use the question generator.")

    messages = _build_messages(payload)
    comp_kwargs = _completion_limit_args(DEFAULT_MODEL)
    try:
        stream = await acompletion(
            model=DEFAULT_MODEL,
            messages=messages,
            temperature=TEMPERATURE,
            stream=True,
            **comp_kwargs,
        )
    except Exception as exc:
        raise QuestionGenerationError(f"LLM request failed: {exc}") from exc

    collected: List[str] = []
    async for chunk in stream:
        delta = chunk["choices"][0]["delta"]
        content = delta.get("content")
        if not content:
            continue
        collected.append(content)
        yield {"type": "chunk", "content": content}

    full_text = "".join(collected)
    questions = _parse_questions(full_text)
    markdown = render_canvas_markdown(payload.instructions, [q.model_dump() for q in questions], {})
    yield {
        "type": "complete",
        "questions": [q.model_dump() for q in questions],
        "markdown": markdown,
        "raw_response": full_text,
    }


def _build_messages(payload: QuestionGenerationRequest) -> List[Dict[str, str]]:
    derived_type_counts, derived_total = _derive_type_counts(payload.instructions)
    type_instruction = "Feel free to use MCQ, short_answer, true_false, or essay question types."
    if payload.question_types:
        type_instruction = (
            f"The instructor prefers the following question types (use these labels when possible): {', '.join(payload.question_types)}."
        )
    elif derived_type_counts:
        summary = ", ".join(f"{count} {kind.replace('_', ' ')}" for kind, count in derived_type_counts.items())
        type_instruction = (
            f"You must generate exactly these question counts/types: {summary}. Use only the labels mcq, true_false, short_answer, or essay."
        )

    total_questions = payload.question_count or derived_total
    count_instruction = "Generate only the requested questions."
    if total_questions is not None:
        count_instruction = f"Generate exactly {total_questions} questions."
    if derived_type_counts and derived_total:
        count_instruction = f"Generate exactly {derived_total} questions total, matching the per-type counts above."

    context_block = f"\nContext:\n{payload.context.strip()}" if payload.context else ""

    schema_block = """
Return JSON with this shape:
{
  "questions": [
    {
      "kind": "mcq | short_answer | true_false | essay",
      "text": "Question text",
      "options": ["optional list of options"],
      "answer": "short answer or letter",
      "explanation": "why the answer is correct",
      "reference": "source citation"
    }
  ]
}
"""

    system_prompt = (
        "You are an experienced instructor who writes exam-ready questions. "
        "Only produce valid JSON. Avoid commentary outside JSON."
    )

    constraint_note = (
        "If you cannot satisfy the requested counts/types, respond with a JSON error object like {\"error\": \"reason\"} instead of returning questions."
    )
    if derived_type_counts:
        constraint_note += " Do not output question types that were not explicitly requested."

    user_prompt = (
        f"{payload.instructions.strip()}\n\n"
        f"{type_instruction}\n"
        f"{count_instruction}\n"
        f"{schema_block}\n"
        f"{context_block}\n"
        f"Ensure answers reflect the context. {constraint_note}"
    )

    return [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_prompt},
    ]


def _parse_questions(content: str) -> List[Question]:
    cleaned = _strip_code_fences(content)
    payload: Any
    try:
        payload = json.loads(cleaned)
    except json.JSONDecodeError as exc:
        raise QuestionGenerationError(f"LLM response was not valid JSON: {exc}") from exc

    raw_questions = payload.get("questions") if isinstance(payload, dict) else payload
    if not isinstance(raw_questions, list):
        raise QuestionGenerationError("LLM response did not include a 'questions' list.")

    normalized: List[Question] = []
    for entry in raw_questions:
        if not isinstance(entry, dict):
            continue
        text = entry.get("text") or entry.get("question")
        if not text:
            continue
        options = entry.get("options")
        if isinstance(options, list):
            options = [str(o).strip() for o in options if str(o).strip()]
        else:
            options = None
        question = Question(
            kind=(entry.get("kind") or entry.get("type") or "short_answer").lower(),
            text=text.strip(),
            options=options,
            answer=(entry.get("answer") or entry.get("solution") or "").strip() or None,
            explanation=(entry.get("explanation") or entry.get("rationale") or "").strip() or None,
            reference=(entry.get("reference") or entry.get("source") or "").strip() or None,
        )
        normalized.append(question)

    if not normalized:
        raise QuestionGenerationError("LLM response did not include any valid questions.")

    return normalized


def _strip_code_fences(raw: str) -> str:
    pattern = re.compile(r"^```(?:json)?\s*(.*?)\s*```$", re.DOTALL)
    match = pattern.match(raw.strip())
    if match:
        return match.group(1)
    return raw


def _derive_type_counts(instructions: str) -> Tuple[Dict[str, int], Optional[int]]:
    counts: Dict[str, int] = {}
    text = instructions.lower()

    for kind, pattern in TYPE_PATTERNS.items():
        regex = re.compile(rf"(\d+)\s*(?:{pattern})", re.I)
        for match in regex.finditer(text):
            qty = int(match.group(1))
            counts[kind] = counts.get(kind, 0) + qty

    total = sum(counts.values()) if counts else None
    if total is None:
        general = re.search(r"(\d+)\s+(?:questions|items)", text)
        if general:
            total = int(general.group(1))

    return counts, total


def summarize_paper_chat(paper_id: int, messages: List[PaperChatMessage]) -> Dict[str, Any]:
    with get_conn() as conn:
        paper = conn.execute("SELECT id, title FROM papers WHERE id=?", (paper_id,)).fetchone()
        if not paper:
            raise QuestionGenerationError("Paper not found.")
        sections = conn.execute(
            "SELECT page_no, text FROM sections WHERE paper_id=? ORDER BY page_no ASC",
            (paper_id,)
        ).fetchall()
    context = "\n\n".join((row["text"] or "" for row in sections))[:MAX_CONTEXT_CHARS]
    if not context.strip():
        raise QuestionGenerationError("No text available for this paper.")

    system_prompt = (
        "You are a research assistant. Summarize the given paper and answer follow-up questions using only the provided context."
    )
    base_messages: List[Dict[str, str]] = [
        {"role": "system", "content": f"{system_prompt}\nContext:\n{context}"}
    ]
    for msg in messages:
        base_messages.append({"role": msg.role, "content": msg.content})

    try:
        response = completion(
            model=DEFAULT_MODEL,
            messages=base_messages,
            temperature=TEMPERATURE,
            **_completion_limit_args(DEFAULT_MODEL),
        )
    except Exception as exc:
        raise QuestionGenerationError(f"LLM request failed: {exc}") from exc

    text = response["choices"][0]["message"]["content"].strip()
    note_title = paper["title"] or "Paper Summary"
    return {
        "message": text,
        "paper_id": paper["id"],
        "paper_title": paper["title"],
        "suggested_title": note_title,
    }


async def extract_context_from_upload(filename: str, data: bytes) -> QuestionContextUploadResponse:
    suffix = Path(filename or "").suffix.lower()
    if suffix not in {".pdf", ".ppt", ".pptx"}:
        raise QuestionGenerationError("Only PDF and PPT/PPTX files are supported at the moment.")

    loop = asyncio.get_running_loop()
    text = await loop.run_in_executor(None, _extract_text_from_bytes, filename, data)
    text = text.strip()
    if not text:
        raise QuestionGenerationError("Could not extract any text from the uploaded file.")
    if len(text) > MAX_CONTEXT_CHARS:
        text = text[:MAX_CONTEXT_CHARS]
    preview = text[:400].strip()
    return QuestionContextUploadResponse(
        context_id=uuid.uuid4().hex,
        filename=filename or "upload",
        characters=len(text),
        preview=preview,
        text=text,
    )


def _extract_text_from_bytes(filename: str, data: bytes) -> str:
    suffix = Path(filename or "").suffix.lower()
    with NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(data)
        tmp_path = Path(tmp.name)
    try:
        if suffix == ".pdf":
            return _extract_pdf_text(tmp_path)
        return _extract_ppt_text(tmp_path)
    finally:
        tmp_path.unlink(missing_ok=True)


def _extract_pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    parts: List[str] = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts)


def _extract_ppt_text(path: Path) -> str:
    presentation = Presentation(str(path))
    parts: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)
